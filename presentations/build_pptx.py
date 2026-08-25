# -*- coding: utf-8 -*-
"""Growlab pitch — شبكة توزيع أدائي (مشاركة + UGC + قصّ) · same CloudMesh theme."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from lxml import etree

# —— Growlab Cloud.Craft (from site screenshot / globals.css) ——
PAPER = RGBColor(0xF7, 0xF8, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUNK = RGBColor(0xEE, 0xF1, 0xF4)
INK = RGBColor(0x11, 0x13, 0x18)
INK_DIM = RGBColor(0x5C, 0x65, 0x73)
INK_FAINT = RGBColor(0x8B, 0x93, 0xA1)
LINE = RGBColor(0xE2, 0xE5, 0xEA)
SIGNAL = RGBColor(0x1F, 0x6F, 0xEB)
SIGNAL_DIM = RGBColor(0xD6, 0xE4, 0xFA)
PULSE = RGBColor(0x6F, 0x9B, 0x7C)
MESH_CYAN = RGBColor(0x7D, 0xD3, 0xFC)
MESH_LIME = RGBColor(0x86, 0xEF, 0xAC)
MESH_SUN = RGBColor(0xFD, 0xE6, 0x8A)
MESH_CORAL = RGBColor(0xFD, 0xBA, 0x74)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
M = Inches(0.5)
CONTENT_W = Inches(12.333)
FONT = "Arial"

# Example performance payouts (illustrative OMR)
EX_VISIT = 0.10
EX_VIEW_CPM = 0.20  # per 1000 views (creator)
EX_BUY_PCT_SHARE = 10
EX_BUY_PCT_ORIGIN = 15
EX_BUY_PCT_CLIP = 10
EX_ORIGIN_BONUS = 3


def pt_emu(pt):
    return int(pt * 12700)


def set_fill_alpha(shape, alpha_pct):
    spPr = shape._element.spPr
    solid = spPr.find(qn("a:solidFill"))
    if solid is None:
        return
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:
        return
    for old in srgb.findall(qn("a:alpha")):
        srgb.remove(old)
    a = etree.SubElement(srgb, qn("a:alpha"))
    a.set("val", str(int(alpha_pct * 1000)))


def soft_blur(shape, soft_pt=90, glow_pt=35, glow_hex="7DD3FC", glow_alpha=35):
    """Diffuse wash — softEdge so no hard circle."""
    spPr = shape._element.spPr
    for old in spPr.findall(qn("a:effectLst")):
        spPr.remove(old)
    effect = etree.SubElement(spPr, qn("a:effectLst"))
    soft = etree.SubElement(effect, qn("a:softEdge"))
    soft.set("rad", str(pt_emu(soft_pt)))
    glow = etree.SubElement(effect, qn("a:glow"))
    glow.set("rad", str(pt_emu(glow_pt)))
    srgb = etree.SubElement(glow, qn("a:srgbClr"))
    srgb.set("val", glow_hex)
    alpha = etree.SubElement(srgb, qn("a:alpha"))
    alpha.set("val", str(int(glow_alpha * 1000)))


def card_shadow(shape):
    """Soft drop shadow like site cards."""
    spPr = shape._element.spPr
    for old in spPr.findall(qn("a:effectLst")):
        spPr.remove(old)
    effect = etree.SubElement(spPr, qn("a:effectLst"))
    sh = etree.SubElement(effect, qn("a:outerShdw"))
    sh.set("blurRad", str(pt_emu(22)))
    sh.set("dist", str(pt_emu(6)))
    sh.set("dir", "2700000")
    sh.set("algn", "ctr")
    sh.set("rotWithShape", "0")
    c = etree.SubElement(sh, qn("a:srgbClr"))
    c.set("val", "0F172A")
    ca = etree.SubElement(c, qn("a:alpha"))
    ca.set("val", "12000")


def send_back(slide, shape, index=3):
    spTree = slide.shapes._spTree
    el = shape._element
    spTree.remove(el)
    spTree.insert(index, el)


def rgb_hex(rgb):
    return f"{int(rgb[0]):02X}{int(rgb[1]):02X}{int(rgb[2]):02X}"


def mesh_wash(slide, left, top, w, h, color, alpha=28, soft=100, glow=40):
    """Pastel CloudMesh orb — blurred glow only."""
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    set_fill_alpha(sh, alpha)
    hx = rgb_hex(color)
    soft_blur(sh, soft_pt=soft, glow_pt=glow, glow_hex=hx, glow_alpha=min(alpha + 10, 45))
    send_back(slide, sh, 3)
    return sh


def set_run(run, size=16, bold=False, color=INK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", FONT)


def rtl(p, align=PP_ALIGN.RIGHT):
    p.alignment = align
    p._p.get_or_add_pPr().set("rtl", "1")


def txt(slide, left, top, w, h, text, size=16, bold=False, color=INK, align=PP_ALIGN.RIGHT):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    rtl(p, align)
    r = p.add_run()
    r.text = text
    set_run(r, size=size, bold=bold, color=color)
    return box


def lines(slide, left, top, w, h, items, size=14, color=INK_DIM, bold=False, after=8):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        rtl(p)
        p.space_after = Pt(after)
        r = p.add_run()
        r.text = item
        set_run(r, size=size, bold=bold, color=color)
    return box


def bg(slide):
    """White paper + soft mesh glows like growlab.om marketing."""
    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    base.fill.solid()
    base.fill.fore_color.rgb = PAPER
    base.line.fill.background()
    send_back(slide, base, 2)

    mesh_wash(slide, Inches(7.5), Inches(-2.2), Inches(8.0), Inches(7.0), MESH_CYAN, alpha=32, soft=110, glow=50)
    mesh_wash(slide, Inches(-3.0), Inches(1.5), Inches(7.0), Inches(6.5), MESH_CORAL, alpha=22, soft=105, glow=45)
    mesh_wash(slide, Inches(2.0), Inches(5.0), Inches(6.5), Inches(4.5), MESH_SUN, alpha=18, soft=100, glow=40)
    mesh_wash(slide, Inches(9.5), Inches(4.5), Inches(5.0), Inches(4.0), MESH_LIME, alpha=16, soft=95, glow=35)


def card(slide, left, top, w, h, fill=WHITE):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = LINE
    sh.line.width = Pt(1)
    try:
        sh.adjustments[0] = 0.1
    except Exception:
        pass
    card_shadow(sh)
    return sh


def accent_bar(slide, left, top, h, color=SIGNAL):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.055), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def eyebrow(slide, text):
    txt(slide, M, Inches(0.38), CONTENT_W, Inches(0.3), text, size=12, color=INK_FAINT)


def title(slide, text, top=Inches(0.68)):
    txt(slide, M, top, CONTENT_W, Inches(0.7), text, size=28, bold=True, color=INK)


def lede(slide, text, top=Inches(1.4)):
    lines(slide, M, top, CONTENT_W, Inches(0.55), [text], size=15, color=INK_DIM, after=2)


def footer(slide, n, total=17):
    txt(slide, M, Inches(7.1), Inches(4), Inches(0.25), "Growlab", size=10, bold=True, color=INK_FAINT)
    txt(slide, Inches(10.5), Inches(7.1), Inches(2.3), Inches(0.25), f"{n} / {total}", size=10, color=INK_FAINT, align=PP_ALIGN.LEFT)


def feature_card(slide, left, top, w, h, num, heading, body, accent=SIGNAL):
    card(slide, left, top, w, h)
    accent_bar(slide, left, top, h, accent)
    y = top + Inches(0.16)
    if num:
        txt(slide, left + Inches(0.22), y, w - Inches(0.35), Inches(0.25), num, size=11, bold=True, color=SIGNAL)
        y += Inches(0.26)
    txt(slide, left + Inches(0.22), y, w - Inches(0.35), Inches(0.35), heading, size=14, bold=True, color=INK)
    y += Inches(0.38)
    lines(slide, left + Inches(0.22), y, w - Inches(0.35), h - (y - top) - Inches(0.12), [body], size=12, color=INK_DIM, after=2)


def style_chart(chart, has_legend=True):
    chart.has_legend = has_legend
    if has_legend:
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(11)
        chart.legend.font.color.rgb = INK_DIM
        chart.legend.font.name = FONT
    plot = chart.plots[0]
    plot.has_data_labels = True
    try:
        dl = plot.data_labels
        dl.font.size = Pt(11)
        dl.font.bold = True
        dl.font.color.rgb = INK
        dl.font.name = FONT
    except Exception:
        pass
    try:
        chart.value_axis.has_major_gridlines = False
        chart.value_axis.tick_labels.font.size = Pt(10)
        chart.value_axis.tick_labels.font.color.rgb = INK_FAINT
        chart.category_axis.tick_labels.font.size = Pt(11)
        chart.category_axis.tick_labels.font.color.rgb = INK_DIM
    except Exception:
        pass


def color_series(chart, colors):
    for i, color in enumerate(colors):
        try:
            series = chart.series[i]
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color
        except Exception:
            pass


def bar_compare(slide, left, top, w, h):
    """Cash out before result vs pay on performance."""
    data = CategoryChartData()
    data.categories = ["كاش إعلان\nمقدماً", "رسوم وكالة\nثابتة", "دفع على\nأداء فقط"]
    data.add_series("الطريقة القديمة", (80, 50, 0))
    data.add_series("Growlab", (0, 0, 35))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, w, h, data).chart
    style_chart(chart, has_legend=True)
    color_series(chart, [INK_FAINT, SIGNAL])
    try:
        chart.plots[0].data_labels.number_format = "0"
    except Exception:
        pass
    return chart


def loop_bars(slide, left, top, w, h):
    """Illustrative growth loop: shares → visits → orders → new sharers."""
    data = CategoryChartData()
    data.categories = ["مشاركات", "زيارات رابط", "مشتريات محصّلة", "مشاركون جدد"]
    data.add_series("العدد (توضيحي)", (200, 95, 28, 22))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, left, top, w, h, data).chart
    style_chart(chart, has_legend=False)
    color_series(chart, [SIGNAL])
    try:
        chart.plots[0].data_labels.number_format = "0"
    except Exception:
        pass
    return chart


def payout_bars(slide, left, top, w, h):
    """Relative earn potential: share vs UGC vs clip (index)."""
    data = CategoryChartData()
    data.categories = ["مشارك\n(رابط)", "صانع\n(UGC)", "قصّاص\n(Clip)"]
    data.add_series("مؤشر فرصة الربح", (70, 100, 85))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, w, h, data).chart
    style_chart(chart, has_legend=False)
    color_series(chart, [SIGNAL])
    try:
        chart.plots[0].data_labels.number_format = "0"
    except Exception:
        pass
    return chart


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    total = 17
    n = 0

    def next_slide():
        nonlocal n
        n += 1
        s = prs.slides.add_slide(blank)
        bg(s)
        return s

    # ── 1 Title ──
    s = next_slide()
    card(s, M, Inches(1.45), Inches(3.4), Inches(0.38), SIGNAL_DIM)
    txt(
        s, M + Inches(0.08), Inches(1.5), Inches(3.25), Inches(0.28),
        "شبكة توزيع أدائي عمانية", size=12, bold=True, color=SIGNAL, align=PP_ALIGN.CENTER,
    )
    txt(s, M, Inches(2.1), Inches(11), Inches(1.0), "Growlab", size=56, bold=True, color=INK)
    txt(s, M, Inches(3.15), Inches(11), Inches(0.55), "عميلك أفضل مسوّق لك.", size=24, bold=True, color=INK_DIM)
    lines(s, M, Inches(3.9), Inches(11), Inches(1.0), [
        "افتح حملة بسقفك. زبون يشتري ويوصّي — وتدفع على زيارة وبيع محصّل فقط. لا إعلان مقدماً ولا وكالة.",
    ], size=15, color=INK_DIM, after=6)
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.55), SLIDE_W, Inches(0.95))
    strip.fill.solid()
    strip.fill.fore_color.rgb = INK
    strip.line.fill.background()
    txt(s, M, Inches(6.75), Inches(9), Inches(0.35), "مشاركة · UGC · قصّ · أداء · دفتر", size=13, color=RGBColor(0xA8, 0x9B, 0x86))
    txt(s, Inches(10), Inches(6.75), Inches(2.8), Inches(0.35), f"١ / {total}", size=12, color=INK_FAINT, align=PP_ALIGN.LEFT)

    # ── 2 الفكرة ──
    s = next_slide()
    eyebrow(s, "الفكرة")
    title(s, "من منتجك إلى شبكة توزيع أدائي")
    lede(s, "أنت في المركز: تفتح الحملة، تحدّد السقف، وتتابع الأداء. الشبكة — مشاركة، محتوى، قصّ — تشتغل لصالحك.")
    roles = [
        ("٠١  التاجر", "حملة + سقف + متجر", "يفتح حملة بسقف، يبني واجهة متجره، ويحدّد تسعير الأداء. يدفع على نتيجة فقط."),
        ("٠٢  زبونك", "يشارك بعد الشراء", "بعد COD محصّل: رابط إحالة فوراً. المشاركة تكفي — التصوير اختياري يضاعف فرصة الربح."),
        ("٠٣  الشبكة", "محتوى وقصّ", "UGC معتمد يدخل مكتبة. القصّاصون يوسّعون الوصول — سلسلة توزيع من كل بيعة."),
    ]
    for i, (num, h, b) in enumerate(roles):
        feature_card(s, M + i * Inches(4.15), Inches(2.2), Inches(4.0), Inches(4.3), num, h, b)
    footer(s, n, total)

    # ── 3 المشكلة ──
    s = next_slide()
    eyebrow(s, "المشكلة")
    title(s, "تدفع مرتين… وما زلت تسأل: من وين الزباين؟")
    lede(s, "واتساب يرتّب الدردشة. الوكالة تبيع وعوداً. «نسوي لك متجر» مستهلك. ما في محرّك جلب زبائن يشتغل لك.")
    problems = [
        ("01", "تمويل إعلان مقدماً", "تدفع كاش قبل ما تعرف إذا الحملة جابت بيعة محصّلة."),
        ("02", "وكالة بلا محاسبة أداء", "وصول وتقارير… والطلبات ما زالت فوضى واتساب."),
        ("03", "«متجر إلكتروني» مستهلك", "ما ناقصك واجهة — ناقصك ناس يوزّعون ودفتر يقفل الفضل."),
        ("04", "زبائنك يضيعون كقناة", "اشترى وراضٍ… وما في نظام يحوّله لشبكة توزيع تعمل لصالحك."),
    ]
    for i, (num, h, b) in enumerate(problems):
        col, row = i % 2, i // 2
        feature_card(s, M + col * Inches(6.2), Inches(2.15) + row * Inches(2.25), Inches(6.0), Inches(2.1), num, h, b)
    footer(s, n, total)

    # ── 4 الحل ──
    s = next_slide()
    eyebrow(s, "الحل")
    title(s, "عميلك أفضل مسوّق — ادفع على أداء بسقف")
    lede(s, "Growlab تضيف قناة توزيع أدائي من زبائنك أنفسهم. واتساب يبقى للتواصل — والمتجر واجهة أنت تبنيها وتتحكم بها.")
    card(s, M, Inches(2.1), Inches(5.5), Inches(4.5))
    lines(s, M + Inches(0.3), Inches(2.3), Inches(5.0), Inches(4.0), [
        "ما يثبّته Growlab:",
        "",
        "• كل مشترٍ يحصل رابط مشاركة بعد التحصيل",
        "• التصوير (UGC) اختياري — يضاعف فرصة الربح",
        "• موافقة على المحتوى قبل دخول المكتبة",
        "• قصّاصون يعيدون توزيع المقاطع المعتمدة",
        "• التاجر يسعّر: زيارة · شراء · مشاهدة",
        "• سقف ميزانية + إيقاف تلقائي",
        "",
        "الزباين يزيدون من السلسلة — مو من فاتورة إعلان.",
    ], size=13, color=INK_DIM, after=3)
    card(s, Inches(6.3), Inches(2.1), Inches(6.5), Inches(4.5))
    txt(s, Inches(6.55), Inches(2.25), Inches(6.0), Inches(0.3), "أين يروح الكاش؟ (توضيحي)", size=13, bold=True, color=INK)
    bar_compare(s, Inches(6.45), Inches(2.55), Inches(6.2), Inches(3.85))
    footer(s, n, total)

    # ── 5 السلسلة اللانهائية (عملية) ──
    s = next_slide()
    eyebrow(s, "محرّك النمو")
    title(s, "سلسلة تتغذى من كل بيعة محصّلة")
    lede(s, "مشاركة → زيارة → شراء → مشارك جديد → محتوى → قصّ → مشاركة أوسع. تتوسع طالما الميزانية والمنتج يستحقان.")
    chain = [
        ("١", "شراء محصّل", "COD يثبت الثقة."),
        ("٢", "رابط مشاركة", "يفتح تلقائياً — بلا تصوير."),
        ("٣", "UGC اختياري", "موافقة → مكتبة."),
        ("٤", "قصّ وتوزيع", "Clippers يوسّعون."),
        ("٥", "مشتري جديد", "الدورة تعيد نفسها."),
    ]
    for i, (num, h, b) in enumerate(chain):
        feature_card(s, M + i * Inches(2.48), Inches(2.25), Inches(2.38), Inches(4.2), num, h, b)
    footer(s, n, total)

    # ── 6 ثلاثة مسارات ربح ──
    s = next_slide()
    eyebrow(s, "مسارات الدخل")
    title(s, "المشاركة للجميع — التصوير يضاعف — القصّ يوسّع")
    lede(s, "ما في أحد مستبعد لأن ما عنده ستوديو. اللي ما يصور يشارك رابط؛ اللي يصور يدخل دور الصانع.")
    paths = [
        ("٠١  مشارك", "رابط فقط", "بعد الشراء: شارك واتساب/ستوري. تربح من زيارة المتجر والشراء المحصّل من رابطك."),
        ("٠٢  صانع Origin", "UGC معتمد", "تصوّر ضمن شروط التاجر، ترفع الرابط للمراجعة، وبعد الموافقة أرباح أعلى + حصة من مبيعات القصّاصين."),
        ("٠٣  قصّاص Clipper", "إعادة توزيع", "تقصّ من المكتبة المعتمدة وتنشر برابطك — فرصة دخل بلا تصوير من الصفر."),
    ]
    for i, (num, h, b) in enumerate(paths):
        feature_card(s, M + i * Inches(4.15), Inches(2.2), Inches(4.0), Inches(4.3), num, h, b, accent=PULSE if i == 1 else SIGNAL)
    footer(s, n, total)

    # ── 7 كيف يعمل ──
    s = next_slide()
    eyebrow(s, "كيف يعمل")
    title(s, "ثلاث خطوات من حملتك إلى مبيعات محصّلة")
    lede(s, "من فتح الحملة، إلى طلبات COD، إلى أداء يُقاس ويُصرف من سقفك.")
    steps = [
        ("٠١", "افتح حملة بسقف", "أضف المنتج، حدّد أسعار الأداء (زيارة/شراء)، وضع سقف ميزانية — بدون اشتراك."),
        ("٠٢", "شبكة التوزيع تشتغل", "زبائنك يشاركون. محتوى UGC معتمد يُقصّ ويُوزَّع. كل طلب COD يصل لوحتك."),
        ("٠٣", "ادفع على الأداء فقط", "بعد تأكيد التحصيل: خصم من سقف الحملة والدفتر. إيقاف تلقائي عند نفاد الميزانية."),
    ]
    for i, (num, h, b) in enumerate(steps):
        feature_card(s, M + i * Inches(4.15), Inches(2.2), Inches(4.0), Inches(4.3), num, h, b)
    footer(s, n, total)

    # ── 8 قمع السلسلة ──
    s = next_slide()
    eyebrow(s, "التتبع")
    title(s, "رابط واحد يقيس المشاركة حتى البيعة")
    lede(s, "مثال توضيحي: المشاركات تولّد زيارات، والزيارات تولّد مشتريات، والمشتريات تولّد مشاركين جدداً.")
    card(s, M, Inches(2.1), Inches(7.2), Inches(4.5))
    txt(s, M + Inches(0.25), Inches(2.25), Inches(6.7), Inches(0.3), "حلقة الجلب (أرقام توضيحية)", size=13, bold=True, color=INK)
    loop_bars(s, M + Inches(0.15), Inches(2.55), Inches(6.9), Inches(3.9))
    kpis = [("٢٠٠", "مشاركة"), ("٩٥", "زيارة"), ("٢٨", "محصّل"), ("٢٢", "مشارك جديد")]
    for i, (v, lab) in enumerate(kpis):
        y = Inches(2.1) + i * Inches(1.15)
        card(s, Inches(9.0), y, Inches(3.8), Inches(1.05))
        accent_bar(s, Inches(9.0), y, Inches(1.05), SIGNAL if i < 3 else PULSE)
        txt(s, Inches(9.25), y + Inches(0.15), Inches(3.4), Inches(0.4), v, size=22, bold=True, color=INK)
        txt(s, Inches(9.25), y + Inches(0.55), Inches(3.4), Inches(0.3), lab, size=12, color=INK_DIM)
    footer(s, n, total)

    # ── 9 مميزات التاجر ──
    s = next_slide()
    eyebrow(s, "للتجار والشركات")
    title(s, "مزايا تشغيلية لصاحب النشاط")
    lede(s, "لا وعود عامة. هذه الآلية التي تحمي محفظتك وتجلب زبائن بأداء.")
    feats = [
        ("سقف ميزانية أنت تضبطه", "حدّد سقف الحملة. الصرف يتوقف تلقائياً عند النفاد — بلا فاتورة مفاجئة."),
        ("دفع على أداء محصّل", "زيارة وبيع بعد تأكيد التحصيل. لا تمويل إعلان مقدماً ولا رسوم ثابتة بلا نتيجة."),
        ("دفتر مكشوف لكل طلب", "مبيعة، أداء الشبكة، حصتك — سطر واحد لكل طلب COD. نفس الأرقام في لوحتك."),
        ("متجرك بخطوات بسيطة", "ابنِ واجهة على growlab.om/m/… — ذكاء اصطناعي، ألوان، عرض، ونشر. COD من المتجر."),
        ("شروط UGC أنت تكتبها", "مدة، إظهار منتج، ممنوعات — ثم موافقة قبل ما يُوزَّع باسم منتجك."),
        ("شبكة توزيع تلقائية", "مشاركة → زيارة → شراء → مشاركة جديدة. ما تحتاج تدير كل مسوّق يدوياً."),
    ]
    for i, (h, b) in enumerate(feats):
        col, row = i % 3, i // 3
        feature_card(s, M + col * Inches(4.15), Inches(2.15) + row * Inches(2.35), Inches(4.0), Inches(2.2), None, h, b)
    footer(s, n, total)

    # ── 10 متجر التاجر ──
    s = next_slide()
    eyebrow(s, "متجرك")
    title(s, "ابنِ متجرك خطوة بخطوة — مثل ووردبريس")
    lede(s, "مثل ووردبريس: اكتب، اختر لوناً، أضف عرضاً، وانشر. الذكاء الاصطناعي يجهّز البداية لك.")
    store_feats = [
        ("٠١  البداية", "ذكاء اصطناعي أو يدوي", "الذكاء الاصطناعي يقترح نصوصاً وألواناً من منتجاتك — يمكنك تعديل كل شيء بعدها."),
        ("٠٢  الهوية", "اسم وقصة المتجر", "ما يراه الزبون أول ما يفتح الرابط. نص عادي بلا HTML."),
        ("٠٣  المظهر", "لون وتخطيط", "اختر لوناً وتخطيطاً — المعاينة تتحدث فوراً."),
        ("٠٤  العرض", "شريط أعلى المتجر", "عروض موسمية أو افتتاح المتجر — أنت تتحكم بالنص."),
        ("٠٥  النشر", "رابط خاص بك", "growlab.om/m/اسم-متجرك — منتجات، COD، وحلقة توزيع أدائي."),
        ("٠٦  التحكم الكامل", "عرضك، أسعارك، وسائطك", "واجهة بسيطة للتاجر. التوزيع الأدائي فوقها — مو بديل عنها."),
    ]
    for i, (num, h, b) in enumerate(store_feats):
        col, row = i % 3, i // 3
        feature_card(s, M + col * Inches(4.15), Inches(2.15) + row * Inches(2.35), Inches(4.0), Inches(2.2), num, h, b, accent=PULSE if i == 0 else SIGNAL)
    footer(s, n, total)

    # ── 11 جدول الأرباح ──
    s = next_slide()
    eyebrow(s, "التسعير")
    title(s, "ادفع على الأداء — بسقف أنت تضبطه")
    lede(s, "بدون اشتراك. بدون رسوم تسجيل. تفتح حملة، تحدّد أسعار الأداء والسقف، وتدفع فقط على زيارة وبيع محصّل.")
    card(s, M, Inches(2.1), Inches(7.0), Inches(4.5))
    # table header
    headers = ["الحدث", "مشارك", "صانع", "قصّاص"]
    widths = [Inches(2.0), Inches(1.5), Inches(1.5), Inches(1.5)]
    x0 = M + Inches(0.25)
    y0 = Inches(2.35)
    for i, (lab, w) in enumerate(zip(headers, widths)):
        txt(s, x0 + sum(widths[:i], Inches(0)), y0, w, Inches(0.35), lab, size=12, bold=True, color=INK)
    rows = [
        ("زيارة متجر من الرابط", f"{EX_VISIT:.2f}", f"{EX_VISIT + 0.02:.2f}", f"{EX_VISIT:.2f}"),
        ("شراء محصّل", f"{EX_BUY_PCT_SHARE}٪", f"{EX_BUY_PCT_ORIGIN}٪", f"{EX_BUY_PCT_CLIP}٪"),
        ("١٠٠٠ مشاهدة مؤهّلة", "—", f"{EX_VIEW_CPM:.2f}", "٠٫١٢"),
        ("Bonus من قصّ المقطع", "—", f"+{EX_ORIGIN_BONUS}٪", "—"),
    ]
    for ri, row in enumerate(rows):
        y = Inches(2.85) + ri * Inches(0.75)
        for ci, (cell, w) in enumerate(zip(row, widths)):
            col = INK if ci == 0 else SIGNAL
            txt(s, x0 + sum(widths[:ci], Inches(0)), y, w, Inches(0.4), cell, size=13, bold=(ci > 0), color=col if ci else INK_DIM)
    card(s, Inches(7.4), Inches(2.1), Inches(5.4), Inches(4.5))
    txt(s, Inches(7.65), Inches(2.25), Inches(5.0), Inches(0.3), "مؤشر فرصة الربح", size=13, bold=True, color=INK)
    payout_bars(s, Inches(7.5), Inches(2.6), Inches(5.1), Inches(3.7))
    footer(s, n, total)

    # ── 12 الحوكمة ──
    s = next_slide()
    eyebrow(s, "الحوكمة والتحكم")
    title(s, "قواعد تحمي محفظة التاجر")
    lede(s, "ست قواعد ثابتة في النظام — ليست وعوداً في العقد.")
    gov = [
        ("سقف ميزانية", "أنت تحدّد سقف الحملة. الصرف يتوقف تلقائياً عند النفاد."),
        ("تسعير أداء مقفول", "أسعار الزيارة والشراء تُقفل عند فتح الحملة. لا تغيير بأثر رجعي."),
        ("دفع بعد التحصيل", "الأداء يُخصم من سقفك فقط بعد تأكيد تحصيل COD. احتياطي إرجاع 14 يوماً."),
        ("دفتر شفاف", "كل طلب سطر مستقل: المبيعة، أداء الشبكة، حصتك. رسوم البوابة صفر على النقد."),
        ("UGC معتمد", "المحتوى يمر بموافقتك قبل النشر. أنت تتحكم بما يُوزَّع باسم منتجك."),
        ("إسناد تلقائي", "كل زيارة وطلب مربوط بالحملة. لا نزاع على الفضل — الدفتر يقفله."),
    ]
    for i, (h, b) in enumerate(gov):
        col, row = i % 3, i // 3
        feature_card(s, M + col * Inches(4.15), Inches(2.15) + row * Inches(2.35), Inches(4.0), Inches(2.2), None, h, b)
    footer(s, n, total)

    # ── 13 مقارنة ──
    s = next_slide()
    eyebrow(s, "المقارنة")
    title(s, "Growlab مقابل الإعلان والوكالة والمتجر")
    lede(s, "Growlab مقابل الخيارات المألوفة — من منظور التاجر.", top=Inches(1.48))
    headers_c = [("المعيار", SUNK), ("مع Growlab", SIGNAL_DIM), ("بدونها", SUNK)]
    widths_c = [Inches(3.2), Inches(4.45), Inches(4.45)]
    x = M
    for (lab, fill), w in zip(headers_c, widths_c):
        card(s, x, Inches(2.35), w, Inches(0.55), fill)
        txt(s, x + Inches(0.15), Inches(2.45), w - Inches(0.25), Inches(0.35), lab, size=13, bold=True, color=INK)
        x += w + Inches(0.12)
    rows_c = [
        ("جلب الزبائن", "سلسلة مشاركة + محتوى + قصّ — زبائنك يصيرون قناة", "إعلان كاش مقدماً أو وكالة بلا ربط بالمبيعات"),
        ("ماذا تشتري", "أداء بسقف ميزانية أنت تضبطه", "وعود وصول أو «متجر جديد»"),
        ("التكلفة", "زيارة وبيع محصّل — إيقاف تلقائي", "رسوم ثابتة حتى لو فشلت الحملة"),
        ("COD والدفتر", "كل طلب سطر مكشوف: مبيعة، أداء، حصتك", "فوضى واتساب وتقارير بلا تحصيل"),
    ]
    for ri, (c, us, them) in enumerate(rows_c):
        y = Inches(3.05) + ri * Inches(0.95)
        vals = [(c, WHITE, True), (us, SIGNAL_DIM, False), (them, SUNK, False)]
        x = M
        for (lab, fill, bold), w in zip(vals, widths_c):
            card(s, x, y, w, Inches(0.85), fill)
            txt(s, x + Inches(0.15), y + Inches(0.22), w - Inches(0.25), Inches(0.45), lab, size=13, bold=bold, color=INK if bold else INK_DIM)
            x += w + Inches(0.12)
    footer(s, n, total)

    # ── 14 خريطة الطريق ──
    s = next_slide()
    eyebrow(s, "خريطة الطريق")
    title(s, "من التسجيل إلى أول حملة توزيع")
    lede(s, "خمس خطوات واضحة للتاجر أو الشركة.")
    card(s, M, Inches(2.15), Inches(6.0), Inches(4.4))
    accent_bar(s, M, Inches(2.15), Inches(4.4), INK)
    txt(s, M + Inches(0.3), Inches(2.35), Inches(5.5), Inches(0.35), "تاجر", size=15, bold=True, color=INK)
    lines(s, M + Inches(0.3), Inches(2.9), Inches(5.5), Inches(3.4), [
        "١  سجّل حساب الشركة",
        "٢  أضف المنتج وافتح حملة",
        "٣  شبكة التوزيع تنطلق",
        "٤  أكّد طلبات COD",
        "٥  راقب الأداء والسقف",
    ], size=14, color=INK_DIM, after=12)
    card(s, Inches(6.85), Inches(2.15), Inches(6.0), Inches(4.4))
    accent_bar(s, Inches(6.85), Inches(2.15), Inches(4.4), SIGNAL)
    txt(s, Inches(7.15), Inches(2.35), Inches(5.5), Inches(0.35), "الشبكة تعمل لك", size=15, bold=True, color=INK)
    lines(s, Inches(7.15), Inches(2.9), Inches(5.5), Inches(3.4), [
        "١  زبائنك يشاركون بعد الشراء",
        "٢  محتوى UGC معتمد يُقصّ ويُوزَّع",
        "٣  كل زيارة وطلب COD يُنسب للحملة",
        "٤  أداء يُخصم من سقفك فقط",
        "٥  حلقة تتكرر بلا إدارة يدوية",
    ], size=14, color=INK_DIM, after=12)
    footer(s, n, total)

    # ── 15 مراحل المنتج ──
    s = next_slide()
    eyebrow(s, "مراحل البناء")
    title(s, "نطلق السلسلة بأمان — ثم نفتح المشاهدات")
    phases = [
        ("v1", "حملة + متجر + مشاركة", "رابط بعد التحصيل، متجر تاجر بسيط، أرباح زيارة وشراء، موافقة UGC، سقف ميزانية."),
        ("v2", "القصّ والمضاعفة", "أدوات قصّ، مكافأة Origin من مبيعات Clipper، شرائح أداء."),
        ("v3", "مشاهدة مؤهّلة", "تحقق أقوى للمشاهدات، اكتشاف مقاطع عالية أداء، حملات فئات."),
    ]
    for i, (num, h, b) in enumerate(phases):
        feature_card(s, M + i * Inches(4.15), Inches(2.2), Inches(4.0), Inches(4.3), num, h, b, accent=SIGNAL if i == 0 else (PULSE if i == 1 else MESH_CORAL))
    footer(s, n, total)

    # ── 16 مؤسس ──
    s = next_slide()
    eyebrow(s, "المؤسس")
    title(s, "مؤسس واحد، بناها شبكة توزيع لا متجر")
    lede(s, "قصي نادر النبراوي — مؤسس Growlab")
    card(s, M, Inches(2.15), Inches(2.7), Inches(0.4), SIGNAL_DIM)
    txt(s, M + Inches(0.1), Inches(2.2), Inches(2.5), Inches(0.3), "طالب أمن سيبراني", size=12, bold=True, color=SIGNAL, align=PP_ALIGN.CENTER)
    card(s, M + Inches(2.9), Inches(2.15), Inches(2.9), Inches(0.4), SUNK)
    txt(s, M + Inches(3.0), Inches(2.2), Inches(2.7), Inches(0.3), "٥ سنوات خبرة تسويق", size=12, bold=True, color=INK_DIM, align=PP_ALIGN.CENTER)
    card(s, M, Inches(2.8), CONTENT_W, Inches(3.6))
    lines(s, M + Inches(0.4), Inches(3.1), Inches(11.5), Inches(2.2), [
        "أنا قصي، طالب أمن سيبراني قبل أي شيء. بنيت Growlab على سؤال «أين ثغرته؟» — لا على شعار تسويقي جاهز.",
        "منصة تدير ميزانيات تجار يجب أن تكون محصّنة من الاحتيال. وخمس سنوات تسويق فعلي تعني نظاماً لتاجر يحتاج زبائن بأداء من اليوم الأول.",
    ], size=17, color=INK_DIM, after=10)
    txt(s, M + Inches(0.4), Inches(5.5), Inches(11.5), Inches(0.5), "«لا نربح إلا إذا ربحت. هذا كل الفرق.»", size=18, bold=True, color=INK)
    footer(s, n, total)

    # ── 17 CTA ──
    s = next_slide()
    eyebrow(s, "لنبدأ")
    title(s, "جاهز تفتح أول حملة؟")
    lede(s, "احجز استشارة مجانية ١٥ دقيقة، أو تواصل معنا مباشرة على واتساب.")
    cta = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, M, Inches(2.3), Inches(5.8), Inches(0.75))
    cta.fill.solid()
    cta.fill.fore_color.rgb = INK
    cta.line.fill.background()
    try:
        cta.adjustments[0] = 0.5
    except Exception:
        pass
    txt(s, M + Inches(0.2), Inches(2.45), Inches(5.4), Inches(0.45), "افتح حساب شركتك", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    card(s, Inches(7.1), Inches(2.3), Inches(5.7), Inches(0.75))
    txt(s, Inches(7.3), Inches(2.45), Inches(5.3), Inches(0.45), "+968 9784 4742", size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
    txt(s, M, Inches(3.25), Inches(8), Inches(0.3), "wa.me/96897844742", size=12, color=INK_FAINT)
    feature_card(s, M, Inches(3.8), Inches(12.0), Inches(2.5), None, "تاجر", "حملة + سقف + متجر + شروط UGC + دفع على أداء محصّل")
    footer(s, n, total)

    out = r"c:\Users\User\Growlab\presentations\Growlab-Pitch.pptx"
    prs.save(out)
    print(out)
    return out


if __name__ == "__main__":
    build()
